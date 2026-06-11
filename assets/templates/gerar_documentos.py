"""
Gerador de templates formais da Dignitatis
Produz: .docx, .odt (via conversão LibreOffice), .pptx
"""

import subprocess
from pathlib import Path
from docx import Document
from docx.shared import Pt, Cm, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Inches, Pt as PPTXPt, Emu
from pptx.dml.color import RGBColor as PPTXRGBColor
from pptx.enum.text import PP_ALIGN

OUT = Path(__file__).parent

# ── Paleta Terracota & Azul-Noite ────────────────────────────────────────────
TERRACOTA   = RGBColor(0xB5, 0x4E, 0x31)   # #B54E31 cor primária
AZUL_NOITE  = RGBColor(0x0F, 0x1F, 0x3D)   # #0F1F3D fundos formais
AREIA       = RGBColor(0xF9, 0xF5, 0xEF)   # #F9F5EF fundo documentos
CINZA_QUENTE= RGBColor(0x9E, 0x8E, 0x7C)   # #9E8E7C captions/rodapé
CARVAO      = RGBColor(0x1C, 0x1C, 0x1C)   # #1C1C1C corpo de texto

# aliases internos (usados nas funções abaixo)
GREEN      = TERRACOTA
DARK_GREEN = AZUL_NOITE
CREAM      = AREIA
GRAY_MID   = CINZA_QUENTE
BLACK_SOFT = CARVAO

PPTX_TERRA  = PPTXRGBColor(0xB5, 0x4E, 0x31)
PPTX_NOITE  = PPTXRGBColor(0x0F, 0x1F, 0x3D)
PPTX_AREIA  = PPTXRGBColor(0xF9, 0xF5, 0xEF)
PPTX_CINZA  = PPTXRGBColor(0x9E, 0x8E, 0x7C)
PPTX_CARVAO = PPTXRGBColor(0x1C, 0x1C, 0x1C)

# aliases para as funções de slides que usam os nomes antigos
PPTX_GREEN      = PPTX_TERRA
PPTX_DARK_GREEN = PPTX_NOITE
PPTX_CREAM      = PPTX_AREIA
PPTX_GRAY       = PPTX_CINZA
PPTX_BLACK      = PPTX_CARVAO

# ── Helpers ──────────────────────────────────────────────────────────────────

def set_cell_bg(cell, hex_color: str):
    """Define cor de fundo de célula de tabela."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), hex_color)
    tcPr.append(shd)


def add_hr(doc: Document, color: str = "B54E31", thickness: int = 12):
    """Adiciona linha horizontal decorativa."""
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(0)
    p.paragraph_format.space_after  = Pt(6)
    pPr = p._p.get_or_add_pPr()
    pBdr = OxmlElement("w:pBdr")
    bottom = OxmlElement("w:bottom")
    bottom.set(qn("w:val"), "single")
    bottom.set(qn("w:sz"), str(thickness))
    bottom.set(qn("w:space"), "1")
    bottom.set(qn("w:color"), color)
    pBdr.append(bottom)
    pPr.append(pBdr)
    return p


def header_block(doc: Document, subtitulo: str = ""):
    """Cabeçalho padrão: nome da organização + subtítulo opcional."""
    # Linha de nome
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("DIGNITATIS")
    run.bold = True
    run.font.size = Pt(18)
    run.font.color.rgb = DARK_GREEN
    run.font.name = "Georgia"

    # Anel multilíngue (tagline)
    p2 = doc.add_paragraph()
    p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run2 = p2.add_run("Direitos Humanos · Paraíba · Brasil")
    run2.font.size = Pt(8)
    run2.font.color.rgb = GRAY_MID
    run2.font.name = "Georgia"

    if subtitulo:
        p3 = doc.add_paragraph()
        p3.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run3 = p3.add_run(subtitulo)
        run3.font.size = Pt(9)
        run3.font.color.rgb = GREEN
        run3.font.name = "Calibri"

    add_hr(doc)


def footer_fields(doc: Document):
    """Rodapé com campos marcadores."""
    add_hr(doc, thickness=6)
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(
        "[Nome Legal da Organização]  ·  CNPJ: [00.000.000/0000-00]  ·  "
        "[Endereço]  ·  [Cidade/UF]  ·  [e-mail]  ·  [site]"
    )
    run.font.size = Pt(7)
    run.font.color.rgb = GRAY_MID
    run.font.name = "Calibri"


def docx_to_odt(docx_path: Path) -> Path:
    """Converte .docx para .odt via LibreOffice headless."""
    odt_path = docx_path.with_suffix(".odt")
    subprocess.run(
        ["libreoffice", "--headless", "--convert-to", "odt",
         "--outdir", str(docx_path.parent), str(docx_path)],
        check=True, capture_output=True
    )
    return odt_path


# ── 1. Ofício ────────────────────────────────────────────────────────────────

def gerar_oficio():
    doc = Document()

    # Margens
    for section in doc.sections:
        section.top_margin    = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin   = Cm(3.0)
        section.right_margin  = Cm(2.0)

    header_block(doc, "Organização de Direitos Humanos")

    # Campos de identificação
    doc.add_paragraph()
    campos = [
        ("Ofício Nº:", "[000/ANO]"),
        ("Data:",      "[Cidade], [DD] de [mês] de [AAAA]"),
        ("Para:",      "[Nome do Destinatário]"),
        ("Cargo:",     "[Cargo/Instituição]"),
        ("Assunto:",   "[Assunto do ofício]"),
    ]
    for label, valor in campos:
        p = doc.add_paragraph()
        r1 = p.add_run(f"{label} ")
        r1.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = DARK_GREEN
        r1.font.name = "Calibri"
        r2 = p.add_run(valor)
        r2.font.size = Pt(10)
        r2.font.color.rgb = BLACK_SOFT
        r2.font.name = "Calibri"

    add_hr(doc, thickness=6)

    # Saudação
    doc.add_paragraph()
    p_sauda = doc.add_paragraph("Prezado(a) Senhor(a),")
    p_sauda.runs[0].font.size = Pt(11)
    p_sauda.runs[0].font.name = "Calibri"

    # Corpo
    doc.add_paragraph()
    for _ in range(3):
        p_corpo = doc.add_paragraph(
            "[Insira aqui o texto do ofício. Este é um parágrafo de exemplo "
            "que deve ser substituído pelo conteúdo real do documento.]"
        )
        p_corpo.runs[0].font.size = Pt(11)
        p_corpo.runs[0].font.color.rgb = BLACK_SOFT
        p_corpo.runs[0].font.name = "Calibri"
        p_corpo.paragraph_format.first_line_indent = Cm(1.25)
        p_corpo.paragraph_format.space_after = Pt(6)

    # Encerramento
    doc.add_paragraph()
    p_enc = doc.add_paragraph("Atenciosamente,")
    p_enc.runs[0].font.size = Pt(11)
    p_enc.runs[0].font.name = "Calibri"

    doc.add_paragraph()
    doc.add_paragraph()
    p_ass = doc.add_paragraph("_____________________________________________")
    p_ass.runs[0].font.color.rgb = DARK_GREEN

    for linha in ["[Nome Completo]", "[Cargo na Organização]", "DIGNITATIS"]:
        pl = doc.add_paragraph(linha)
        pl.runs[0].font.size = Pt(10)
        pl.runs[0].font.name = "Calibri"
        if linha == "DIGNITATIS":
            pl.runs[0].bold = True
            pl.runs[0].font.color.rgb = GREEN

    footer_fields(doc)

    path = OUT / "dignitatis_oficio.docx"
    doc.save(path)
    print(f"  ✓ {path.name}")
    return path


# ── 2. Relatório ─────────────────────────────────────────────────────────────

def gerar_relatorio():
    doc = Document()

    for section in doc.sections:
        section.top_margin    = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin   = Cm(3.0)
        section.right_margin  = Cm(2.0)

    # Capa
    doc.add_paragraph()
    doc.add_paragraph()
    p_tipo = doc.add_paragraph("RELATÓRIO")
    p_tipo.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p_tipo.runs[0]
    r.bold = True; r.font.size = Pt(28); r.font.color.rgb = GREEN; r.font.name = "Georgia"

    p_titulo = doc.add_paragraph("[Título do Relatório]")
    p_titulo.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = p_titulo.runs[0]
    r2.font.size = Pt(18); r2.font.color.rgb = DARK_GREEN; r2.font.name = "Georgia"

    doc.add_paragraph()
    p_sub = doc.add_paragraph("Direitos Humanos · Paraíba · Brasil")
    p_sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p_sub.runs[0].font.size = Pt(10); p_sub.runs[0].font.color.rgb = GRAY_MID
    p_sub.runs[0].font.name = "Georgia"

    doc.add_paragraph()
    doc.add_paragraph()
    for meta in ["Período: [mês/AAAA – mês/AAAA]", "Elaborado por: [Nome/Equipe]", "Data: [DD/MM/AAAA]"]:
        pm = doc.add_paragraph(meta)
        pm.alignment = WD_ALIGN_PARAGRAPH.CENTER
        pm.runs[0].font.size = Pt(11); pm.runs[0].font.name = "Calibri"
        pm.runs[0].font.color.rgb = BLACK_SOFT

    add_hr(doc)

    p_org = doc.add_paragraph("DIGNITATIS — Organização de Direitos Humanos")
    p_org.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p_org.runs[0].bold = True; p_org.runs[0].font.size = Pt(10)
    p_org.runs[0].font.color.rgb = DARK_GREEN; p_org.runs[0].font.name = "Georgia"

    doc.add_page_break()

    # Sumário
    header_block(doc, "Organização de Direitos Humanos")
    p_sum = doc.add_paragraph("SUMÁRIO")
    p_sum.runs[0].bold = True; p_sum.runs[0].font.size = Pt(14)
    p_sum.runs[0].font.color.rgb = GREEN; p_sum.runs[0].font.name = "Georgia"

    for i, secao in enumerate([
        "Apresentação", "Contexto e Objetivos", "Metodologia",
        "Resultados e Análise", "Conclusões e Recomendações",
        "Referências", "Anexos"
    ], 1):
        ps = doc.add_paragraph(f"{i}. {secao}")
        ps.runs[0].font.size = Pt(11); ps.runs[0].font.name = "Calibri"
        ps.runs[0].font.color.rgb = BLACK_SOFT

    doc.add_page_break()

    # Seções de conteúdo
    for titulo_sec in ["1. Apresentação", "2. Contexto e Objetivos", "3. Metodologia"]:
        header_block(doc, "Organização de Direitos Humanos")
        p_sec = doc.add_paragraph(titulo_sec)
        p_sec.runs[0].bold = True; p_sec.runs[0].font.size = Pt(14)
        p_sec.runs[0].font.color.rgb = GREEN; p_sec.runs[0].font.name = "Georgia"

        for _ in range(2):
            pc = doc.add_paragraph(
                "[Insira aqui o conteúdo desta seção. Substitua este texto pelo "
                "conteúdo real do relatório. Você pode adicionar tabelas, gráficos "
                "e outros elementos conforme necessário.]"
            )
            pc.runs[0].font.size = Pt(11); pc.runs[0].font.name = "Calibri"
            pc.runs[0].font.color.rgb = BLACK_SOFT
            pc.paragraph_format.first_line_indent = Cm(1.25)
            pc.paragraph_format.space_after = Pt(6)

        footer_fields(doc)
        doc.add_page_break()

    path = OUT / "dignitatis_relatorio.docx"
    doc.save(path)
    print(f"  ✓ {path.name}")
    return path


# ── 3. Comunicado / Carta ────────────────────────────────────────────────────

def gerar_comunicado():
    doc = Document()

    for section in doc.sections:
        section.top_margin    = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin   = Cm(3.0)
        section.right_margin  = Cm(2.5)

    header_block(doc, "Comunicado Institucional")

    doc.add_paragraph()
    p_local = doc.add_paragraph("[Cidade], [DD] de [mês] de [AAAA]")
    p_local.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    p_local.runs[0].font.size = Pt(10); p_local.runs[0].font.name = "Calibri"

    doc.add_paragraph()
    p_dest = doc.add_paragraph("A quem possa interessar,")
    p_dest.runs[0].font.size = Pt(11); p_dest.runs[0].font.name = "Calibri"

    doc.add_paragraph()
    for _ in range(4):
        pc = doc.add_paragraph(
            "[Insira aqui o texto do comunicado. Este parágrafo deve ser "
            "substituído pelo conteúdo real.]"
        )
        pc.runs[0].font.size = Pt(11); pc.runs[0].font.name = "Calibri"
        pc.runs[0].font.color.rgb = BLACK_SOFT
        pc.paragraph_format.first_line_indent = Cm(1.25)
        pc.paragraph_format.space_after = Pt(6)

    doc.add_paragraph()
    p_fecho = doc.add_paragraph("Respeitosamente,")
    p_fecho.runs[0].font.size = Pt(11); p_fecho.runs[0].font.name = "Calibri"

    doc.add_paragraph()
    doc.add_paragraph()
    for linha in ["[Nome Completo]", "[Cargo]", "DIGNITATIS"]:
        pl = doc.add_paragraph(linha)
        pl.runs[0].font.size = Pt(10); pl.runs[0].font.name = "Calibri"
        if linha == "DIGNITATIS":
            pl.runs[0].bold = True; pl.runs[0].font.color.rgb = GREEN

    footer_fields(doc)

    path = OUT / "dignitatis_comunicado.docx"
    doc.save(path)
    print(f"  ✓ {path.name}")
    return path


# ── 4. Apresentação PPTX ─────────────────────────────────────────────────────

def hex_to_pptx(h: str) -> PPTXRGBColor:
    h = h.lstrip("#")
    return PPTXRGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_pptx_bg(slide, color: PPTXRGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, text, left, top, width, height,
                 font_name="Fraunces", font_size=24, bold=False,
                 color: PPTXRGBColor = None, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = PPTXPt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox

def add_rect(slide, left, top, width, height, fill_color: PPTXRGBColor, line=False):
    from pptx.util import Pt as PPt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not line:
        shape.line.fill.background()
    return shape

def gerar_apresentacao():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    W = prs.slide_width
    H = prs.slide_height
    blank = prs.slide_layouts[6]  # layout em branco

    # ── Slide 1: Capa ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_pptx_bg(slide, PPTX_DARK_GREEN)

    # Barra lateral verde clara
    add_rect(slide, Inches(0), Inches(0), Inches(0.25), H, PPTX_GREEN)

    # Nome da org
    add_text_box(
        slide, "DIGNITATIS",
        Inches(1), Inches(1.5), Inches(11), Inches(1.2),
        font_name="Fraunces", font_size=54, bold=True,
        color=PPTX_CREAM, align=PP_ALIGN.LEFT
    )
    # Tagline
    add_text_box(
        slide, "Direitos Humanos · Paraíba · Brasil",
        Inches(1), Inches(2.9), Inches(11), Inches(0.5),
        font_name="Fraunces", font_size=13,
        color=PPTX_GRAY, align=PP_ALIGN.LEFT
    )
    # Linha verde
    add_rect(slide, Inches(1), Inches(3.6), Inches(4), Inches(0.04), PPTX_GREEN)

    # Título da apresentação
    add_text_box(
        slide, "[Título da Apresentação]",
        Inches(1), Inches(3.8), Inches(11), Inches(1),
        font_name="Source Sans Pro", font_size=22,
        color=PPTX_CREAM, align=PP_ALIGN.LEFT
    )
    add_text_box(
        slide, "[Subtítulo / Evento / Data]",
        Inches(1), Inches(4.7), Inches(11), Inches(0.7),
        font_name="Source Sans Pro", font_size=14,
        color=PPTX_GRAY, align=PP_ALIGN.LEFT
    )

    # ── Slide 2: Sumário ─────────────────────────────────────────────────────
    slide2 = prs.slides.add_slide(blank)
    add_pptx_bg(slide2, PPTX_CREAM)
    add_rect(slide2, Inches(0), Inches(0), Inches(0.25), H, PPTX_GREEN)

    add_text_box(
        slide2, "SUMÁRIO",
        Inches(1), Inches(0.4), Inches(11), Inches(0.8),
        font_name="Fraunces", font_size=28, bold=True,
        color=PPTX_DARK_GREEN, align=PP_ALIGN.LEFT
    )
    add_rect(slide2, Inches(1), Inches(1.2), Inches(5), Inches(0.04), PPTX_GREEN)

    topicos = [
        "01  Apresentação da Organização",
        "02  Contexto e Problema",
        "03  Nossa Atuação",
        "04  Resultados e Impacto",
        "05  Próximos Passos",
    ]
    for i, topico in enumerate(topicos):
        add_text_box(
            slide2, topico,
            Inches(1), Inches(1.5 + i * 0.85), Inches(11), Inches(0.7),
            font_name="Source Sans Pro", font_size=16,
            color=PPTX_BLACK, align=PP_ALIGN.LEFT
        )

    # ── Slide 3: Conteúdo padrão ─────────────────────────────────────────────
    slide3 = prs.slides.add_slide(blank)
    add_pptx_bg(slide3, PPTX_CREAM)
    add_rect(slide3, Inches(0), Inches(0), Inches(0.25), H, PPTX_GREEN)

    add_text_box(
        slide3, "[Título da Seção]",
        Inches(1), Inches(0.4), Inches(11), Inches(0.8),
        font_name="Fraunces", font_size=28, bold=True,
        color=PPTX_DARK_GREEN, align=PP_ALIGN.LEFT
    )
    add_rect(slide3, Inches(1), Inches(1.2), Inches(5), Inches(0.04), PPTX_GREEN)

    add_text_box(
        slide3,
        "• [Ponto principal 1]\n• [Ponto principal 2]\n• [Ponto principal 3]\n• [Ponto principal 4]",
        Inches(1), Inches(1.4), Inches(6), Inches(4),
        font_name="Source Sans Pro", font_size=16,
        color=PPTX_BLACK, align=PP_ALIGN.LEFT
    )

    # Caixa de destaque
    add_rect(slide3, Inches(8), Inches(1.4), Inches(4.5), Inches(4), PPTX_GREEN)
    add_text_box(
        slide3, "[Dado ou citação em destaque]",
        Inches(8.2), Inches(2.2), Inches(4.1), Inches(2.5),
        font_name="Fraunces", font_size=18,
        color=PPTX_CREAM, align=PP_ALIGN.CENTER
    )

    # ── Slide 4: Citação / Destaque ──────────────────────────────────────────
    slide4 = prs.slides.add_slide(blank)
    add_pptx_bg(slide4, PPTX_TERRA)
    add_rect(slide4, Inches(0), Inches(0), Inches(0.25), H, PPTX_NOITE)

    add_text_box(
        slide4, "“",
        Inches(1), Inches(0.8), Inches(2), Inches(1.5),
        font_name="Fraunces", font_size=96, bold=True,
        color=PPTX_CREAM, align=PP_ALIGN.LEFT
    )
    add_text_box(
        slide4, "[Insira aqui uma citação impactante ou dado relevante sobre direitos humanos.]",
        Inches(1.5), Inches(1.8), Inches(10), Inches(2.5),
        font_name="Fraunces", font_size=22,
        color=PPTX_CREAM, align=PP_ALIGN.LEFT
    )
    add_text_box(
        slide4, "— [Fonte / Autoria]",
        Inches(1.5), Inches(4.5), Inches(10), Inches(0.6),
        font_name="Source Sans Pro", font_size=13,
        color=hex_to_pptx("F9F5EF"), align=PP_ALIGN.LEFT
    )

    # ── Slide 5: Encerramento ────────────────────────────────────────────────
    slide5 = prs.slides.add_slide(blank)
    add_pptx_bg(slide5, PPTX_DARK_GREEN)
    add_rect(slide5, Inches(0), Inches(0), Inches(0.25), H, PPTX_GREEN)

    add_text_box(
        slide5, "Obrigado(a)",
        Inches(1), Inches(1.8), Inches(11), Inches(1.2),
        font_name="Fraunces", font_size=48, bold=True,
        color=PPTX_CREAM, align=PP_ALIGN.LEFT
    )
    add_rect(slide5, Inches(1), Inches(3.1), Inches(3), Inches(0.04), PPTX_GREEN)

    add_text_box(
        slide5, "DIGNITATIS",
        Inches(1), Inches(3.3), Inches(11), Inches(0.7),
        font_name="Fraunces", font_size=18, bold=True,
        color=PPTX_GREEN, align=PP_ALIGN.LEFT
    )
    add_text_box(
        slide5, "Direitos Humanos · Paraíba · Brasil",
        Inches(1), Inches(4.0), Inches(11), Inches(0.5),
        font_name="Fraunces", font_size=11,
        color=PPTX_GRAY, align=PP_ALIGN.LEFT
    )
    add_text_box(
        slide5, "[e-mail]  ·  [site]  ·  [redes sociais]",
        Inches(1), Inches(4.7), Inches(11), Inches(0.5),
        font_name="Source Sans Pro", font_size=12,
        color=PPTX_GRAY, align=PP_ALIGN.LEFT
    )

    path = OUT / "dignitatis_apresentacao.pptx"
    prs.save(path)
    print(f"  ✓ {path.name}")
    return path


# ── Main ─────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    print("Gerando templates Dignitatis...")

    paths_docx = []
    for fn in [gerar_oficio, gerar_relatorio, gerar_comunicado]:
        p = fn()
        paths_docx.append(p)

    gerar_apresentacao()

    print("\nConvertendo para ODT via LibreOffice...")
    for p in paths_docx:
        try:
            odt = docx_to_odt(p)
            print(f"  ✓ {odt.name}")
        except Exception as e:
            print(f"  ✗ Erro ao converter {p.name}: {e}")

    print("\nTodos os templates gerados em:", OUT)
